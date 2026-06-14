"""
build_pptx.py — Generate "Construction Management Process" PPTX from scratch.

Usage:
    pip install python-pptx
    python build_pptx.py

Output: "Construction Management Process.pptx" in the same folder.

Strategy:
- 14 slides mirroring cons-mgmt-process.html.
- Heavy on shapes/diagrams, light on text.
- Each slide: header (number + eyebrow + rule), title, big diagram, caption strip.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ─── design tokens (resemble OKLCH → RGB approximations) ───────────────────
INK       = RGBColor(0x2E, 0x33, 0x42)
INK2      = RGBColor(0x53, 0x5A, 0x6B)
MUTED     = RGBColor(0x78, 0x80, 0x91)
FAINT     = RGBColor(0xA8, 0xAE, 0xBA)
LINE      = RGBColor(0xE5, 0xE7, 0xEB)
SURF      = RGBColor(0xFF, 0xFF, 0xFF)
SURF2     = RGBColor(0xF6, 0xF7, 0xF9)

BLUE      = RGBColor(0x4A, 0x6E, 0xCD)
BLUE_SFT  = RGBColor(0xE3, 0xEA, 0xF8)
BLUE_INK  = RGBColor(0x2E, 0x4E, 0xA4)
TEAL      = RGBColor(0x42, 0x9B, 0xAD)
TEAL_SFT  = RGBColor(0xE0, 0xF2, 0xF5)
TEAL_INK  = RGBColor(0x2A, 0x6E, 0x7C)
AMBER     = RGBColor(0xD8, 0x96, 0x33)
AMBER_SFT = RGBColor(0xFA, 0xEF, 0xDC)
AMBER_INK = RGBColor(0x88, 0x59, 0x0F)
ROSE      = RGBColor(0xC2, 0x4F, 0x44)
ROSE_SFT  = RGBColor(0xFA, 0xE3, 0xDF)
ROSE_INK  = RGBColor(0x8F, 0x2E, 0x26)
GREEN     = RGBColor(0x4F, 0xA1, 0x65)
GREEN_SFT = RGBColor(0xE4, 0xF1, 0xE7)
GREEN_INK = RGBColor(0x2E, 0x70, 0x42)
VIOLET    = RGBColor(0x7A, 0x65, 0xC8)
VIOLET_SFT= RGBColor(0xEB, 0xE7, 0xF9)
VIOLET_INK= RGBColor(0x4E, 0x39, 0x9F)

WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)

SANS = "Calibri"
MONO = "Consolas"


# ─── helpers ───────────────────────────────────────────────────────────────
def set_fill(shape, rgb):
    f = shape.fill
    f.solid()
    f.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    ln = shape.line
    ln.color.rgb = rgb
    ln.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_text(shape, text, *, font=SANS, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""  # clear
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kw):
    tb = slide.shapes.add_textbox(left, top, width, height)
    add_text(tb, text, **kw)
    return tb


def add_rect(slide, left, top, width, height, *, fill=SURF, line=LINE, line_w=1.0, radius=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    set_fill(sh, fill)
    if line is None:
        no_line(sh)
    else:
        set_line(sh, line, line_w)
    sh.text_frame.text = ""
    return sh


def add_arrow(slide, x1, y1, x2, y2, color=INK2, width=1.5):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    # add arrowhead at end via XML
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


# ─── slide header helpers ──────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.5)


def slide_header(slide, num, eyebrow, total=14, cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK):
    """Number chip + eyebrow + rule + progress."""
    # number chip
    chip = add_rect(slide, MARGIN, Inches(0.4), Inches(0.6), Inches(0.3),
                    fill=cat_soft, line=None, radius=True)
    add_text(chip, f"{num:02d}", font=MONO, size=11, bold=True, color=cat_ink,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # eyebrow
    eb = add_textbox(slide, MARGIN + Inches(0.75), Inches(0.4), Inches(4.5), Inches(0.3),
                     eyebrow.upper(), font=MONO, size=10, bold=True, color=cat_color)

    # progress on right
    pg = add_textbox(slide, Inches(12.0), Inches(0.4), Inches(0.9), Inches(0.3),
                     f"{num:02d} · {total}", font=MONO, size=10, color=FAINT,
                     align=PP_ALIGN.RIGHT)

    # rule
    rule = slide.shapes.add_connector(1,
                                       MARGIN + Inches(5.4), Inches(0.55),
                                       Inches(11.9), Inches(0.55))
    rule.line.color.rgb = LINE
    rule.line.width = Pt(0.5)


def slide_title(slide, title, sub=None, top=Inches(0.9)):
    add_textbox(slide, MARGIN, top, Inches(12.3), Inches(0.7),
                title, font=SANS, size=28, bold=True, color=INK)
    if sub:
        add_textbox(slide, MARGIN, top + Inches(0.75), Inches(11.5), Inches(0.5),
                    sub, font=SANS, size=13, color=MUTED)


def slide_caption(slide, tag, caption, cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK):
    top = Inches(6.7)
    rule = slide.shapes.add_connector(1, MARGIN, top, Inches(12.83), top)
    rule.line.color.rgb = LINE
    rule.line.width = Pt(0.5)

    tag_box = add_rect(slide, MARGIN, top + Inches(0.15), Inches(1.0), Inches(0.32),
                       fill=cat_soft, line=None, radius=True)
    add_text(tag_box, tag.upper(), font=MONO, size=9, bold=True, color=cat_ink,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, MARGIN + Inches(1.15), top + Inches(0.12), Inches(11.5), Inches(0.55),
                caption, font=SANS, size=12, color=INK2)


# ─── slide builders ────────────────────────────────────────────────────────
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_header(s, 1, "Pengantar", cat_color=INK, cat_soft=SURF2, cat_ink=INK)

    # decorative eyebrow with rules
    add_textbox(s, Inches(0), Inches(2.4), SLIDE_W, Inches(0.4),
                "—  5D BIM · CONSTRUCTION MANAGEMENT PROCESS  —",
                font=MONO, size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # title
    add_textbox(s, Inches(0), Inches(2.9), SLIDE_W, Inches(1.5),
                "Dari Model 3D ke Kurva S",
                font=SANS, size=56, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # subtitle
    add_textbox(s, Inches(2), Inches(4.4), Inches(9.33), Inches(0.6),
                "Empat tahap utama membuat estimasi proyek konstruksi: Harga Satuan → RAB → Jadwal → Kurva S.",
                font=SANS, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2), Inches(4.8), Inches(9.33), Inches(0.6),
                "Setiap tahap diilustrasikan dari nol.",
                font=SANS, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # meta row
    add_textbox(s, Inches(0), Inches(5.6), SLIDE_W, Inches(0.4),
                "📐  Mengajarkan dari nol     🔄  Alur proses lengkap     📊  Banyak diagram, sedikit kata",
                font=SANS, size=12, color=INK, align=PP_ALIGN.CENTER)

    slide_caption(s, "Mulai",
                  "Setiap slide menjelaskan satu tahap dengan diagram dan contoh konkret.",
                  cat_color=INK, cat_soft=SURF2, cat_ink=INK)


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 2, "Big Picture", cat_color=VIOLET, cat_soft=VIOLET_SFT, cat_ink=VIOLET_INK)
    slide_title(s, "5 Tahapan dari Gambar 3D Menjadi Estimasi Lengkap",
                "Setiap tahap memakai output dari tahap sebelumnya. Ini peta besar yang akan kita pelajari satu per satu.")

    # 5 boxes
    boxes = [
        ("3D BIM · REVIT",   "Volume",        "m³, kg, m²",            BLUE_SFT,  BLUE,  BLUE_INK),
        ("TAHAP 1 · AHSP",   "Harga Satuan",  "Rp / m³, Rp / kg",      BLUE_SFT,  BLUE,  BLUE_INK),
        ("TAHAP 2 · RAB",    "Total Biaya",   "Rp total",              TEAL_SFT,  TEAL,  TEAL_INK),
        ("TAHAP 3 · JADWAL", "Gantt",         "minggu × WBS",          AMBER_SFT, AMBER, AMBER_INK),
        ("TAHAP 4 · KURVA S","Progres %",     "akumulasi mingguan",    ROSE_SFT,  ROSE,  ROSE_INK),
    ]

    bw = Inches(2.2); bh = Inches(1.2)
    gap = Inches(0.25)
    total = bw * 5 + gap * 4
    start = (SLIDE_W - total) / 2
    top = Inches(3.0)

    centers = []
    for i, (eyebrow, title, sub, fill, line, ink) in enumerate(boxes):
        left = start + (bw + gap) * i
        box = add_rect(s, left, top, bw, bh, fill=fill, line=line, line_w=1.5, radius=True)
        # no text in box rectangle; layer textboxes for control
        add_textbox(s, left, top + Inches(0.12), bw, Inches(0.3),
                    eyebrow, font=MONO, size=9, bold=True, color=ink, align=PP_ALIGN.CENTER)
        add_textbox(s, left, top + Inches(0.45), bw, Inches(0.4),
                    title, font=SANS, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(s, left, top + Inches(0.85), bw, Inches(0.3),
                    sub, font=MONO, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        centers.append((left, top, bw, bh))

    # arrows between
    for i in range(4):
        x1 = centers[i][0] + bw
        x2 = centers[i + 1][0]
        ym = top + bh / 2
        add_arrow(s, x1, ym, x2, ym, color=VIOLET, width=2)

    # bottom descriptors
    descriptors = ["Hasil QTO", "Σ koef × harga", "Volume × Harga", "Volume ÷ produktivitas", "Σ bobot/minggu"]
    for i, txt in enumerate(descriptors):
        left = start + (bw + gap) * i
        add_textbox(s, left, top + bh + Inches(0.2), bw, Inches(0.3),
                    txt, font=MONO, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    slide_caption(s, "Roadmap",
                  "Volume dari Revit BIM adalah modal awal. Mengubah satu tahap = semua hilirnya ikut berubah.",
                  cat_color=VIOLET, cat_soft=VIOLET_SFT, cat_ink=VIOLET_INK)


def slide_ahsp_sources(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 3, "Tahap 1 · Harga Satuan", cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK)
    slide_title(s, "Dari Mana Harga Bahan & Upah Diambil?",
                "Tiga sumber resmi yang dipakai estimator untuk menyusun Analisa Harga Satuan Pekerjaan (AHSP).")

    # 3 source boxes left
    sources = [
        ("SUMBER 1", "HSPK Pemerintah Daerah", "Harga Satuan Pokok Kegiatan — dirilis tiap awal tahun"),
        ("SUMBER 2", "Survei Pasar (min. 3 toko)", "Harga aktual + ongkos angkut ke lokasi"),
        ("SUMBER 3", "SNI 7394 / 6897 / 2837", "Koefisien resmi: bahan & upah per 1 satuan"),
    ]
    top0 = Inches(3.0); bh = Inches(0.95); gap = Inches(0.15)
    bw = Inches(4.2); left = Inches(0.6)
    centers = []
    for i, (eb, name, desc) in enumerate(sources):
        tp = top0 + (bh + gap) * i
        add_rect(s, left, tp, bw, bh, fill=BLUE_SFT, line=BLUE, line_w=1.5, radius=True)
        add_textbox(s, left + Inches(0.2), tp + Inches(0.08), bw - Inches(0.4), Inches(0.25),
                    eb, font=MONO, size=10, bold=True, color=BLUE_INK)
        add_textbox(s, left + Inches(0.2), tp + Inches(0.32), bw - Inches(0.4), Inches(0.3),
                    name, font=SANS, size=14, bold=True, color=INK)
        add_textbox(s, left + Inches(0.2), tp + Inches(0.62), bw - Inches(0.4), Inches(0.3),
                    desc, font=SANS, size=11, color=MUTED)
        centers.append((left + bw, tp + bh / 2))

    # center AHSP node
    cx = Inches(7.0); cy = Inches(3.5); cw = Inches(3.5); ch = Inches(2.5)
    add_rect(s, cx, cy, cw, ch, fill=BLUE, line=None, radius=True)
    add_textbox(s, cx, cy + Inches(0.4), cw, Inches(0.3),
                "DIPROSES MENJADI", font=MONO, size=10, bold=True, color=BLUE_SFT, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(0.9), cw, Inches(0.6),
                "AHSP", font=SANS, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(1.6), cw, Inches(0.4),
                "Analisa Harga Satuan Pekerjaan",
                font=SANS, size=12, color=BLUE_SFT, align=PP_ALIGN.CENTER)

    # arrows from sources to center
    for (sx, sy) in centers:
        add_arrow(s, sx, sy, cx, cy + ch / 2, color=BLUE, width=1.5)

    # output arrow + label
    add_arrow(s, cx + cw, cy + ch / 2, cx + cw + Inches(0.8), cy + ch / 2, color=BLUE, width=2)
    add_textbox(s, cx + cw + Inches(0.85), cy + ch / 2 - Inches(0.4), Inches(2.0), Inches(0.3),
                "HASIL", font=MONO, size=10, bold=True, color=BLUE_INK)
    add_textbox(s, cx + cw + Inches(0.85), cy + ch / 2 - Inches(0.1), Inches(2.0), Inches(0.4),
                "Rp / m³", font=SANS, size=18, bold=True, color=INK)
    add_textbox(s, cx + cw + Inches(0.85), cy + ch / 2 + Inches(0.35), Inches(2.0), Inches(0.3),
                "per 1 satuan", font=SANS, size=11, color=MUTED)

    slide_caption(s, "Sumber Data",
                  "Proyek pemerintah wajib pakai HSPK. Swasta sering menggabungkan HSPK + survei sendiri + margin negosiasi.",
                  cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK)


def slide_ahsp_formula(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 4, "Tahap 1 · Harga Satuan", cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK)
    slide_title(s, "Cara Menyusun Harga Satuan: Resep yang Kuantitatif",
                "Bayangkan resep masakan: bahan + cara + upah juru masak. AHSP adalah versi terstruktur dari resep ini.")

    # 3 stacked component boxes (left side)
    comps = [
        ("A", "Bahan (Material)",        "Σ koefisien × harga bahan",       BLUE_SFT,  BLUE),
        ("B", "Upah (Labour)",           "Σ koefisien OH × tarif/hari",     TEAL_SFT,  TEAL),
        ("C", "Alat (Equipment, opsional)", "Σ koefisien × sewa alat/hari",  AMBER_SFT, AMBER),
    ]
    left = Inches(0.6); top0 = Inches(2.9); bw = Inches(4.5); bh = Inches(0.85); gap = Inches(0.18)
    for i, (letter, name, sub, fill, line) in enumerate(comps):
        tp = top0 + (bh + gap) * i
        add_rect(s, left, tp, bw, bh, fill=fill, line=line, line_w=1.2, radius=True)
        # circle badge
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.18), tp + Inches(0.15),
                                    Inches(0.55), Inches(0.55))
        set_fill(circle, line); no_line(circle)
        add_text(circle, letter, font=SANS, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, left + Inches(0.85), tp + Inches(0.12), bw - Inches(1), Inches(0.35),
                    name, font=SANS, size=13, bold=True, color=INK)
        add_textbox(s, left + Inches(0.85), tp + Inches(0.45), bw - Inches(1), Inches(0.3),
                    sub, font=MONO, size=10, color=MUTED)

    # plus signs between
    for i in range(2):
        tp = top0 + bh + gap / 2 + (bh + gap) * i - Inches(0.05)
        add_textbox(s, left + bw / 2 - Inches(0.2), tp, Inches(0.4), Inches(0.3),
                    "+", font=MONO, size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # formula text on right
    fx = Inches(5.6); fy = Inches(3.0)
    add_textbox(s, fx, fy, Inches(7), Inches(0.5),
                "Subtotal langsung × (1 + 10% OHP)", font=SANS, size=18, bold=True, color=INK)
    add_textbox(s, fx, fy + Inches(0.55), Inches(7), Inches(0.4),
                "OHP = Overhead & Profit kontraktor (sudah disisipkan oleh AHSP)",
                font=MONO, size=11, color=MUTED)

    # result box
    rx = Inches(8.5); ry = Inches(4.2); rw = Inches(3.8); rh = Inches(1.7)
    add_rect(s, rx, ry, rw, rh, fill=BLUE, line=None, radius=True)
    add_textbox(s, rx, ry + Inches(0.2), rw, Inches(0.3),
                "HARGA SATUAN", font=MONO, size=10, bold=True, color=BLUE_SFT, align=PP_ALIGN.CENTER)
    add_textbox(s, rx, ry + Inches(0.55), rw, Inches(0.6),
                "Rp / m³", font=SANS, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, rx, ry + Inches(1.2), rw, Inches(0.3),
                "≈ Rp 1.320.259 untuk 1 m³ beton K-225",
                font=MONO, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # footnote "apa itu koefisien?"
    add_textbox(s, Inches(5.6), Inches(4.2), Inches(2.8), Inches(0.3),
                "APA ITU KOEFISIEN?", font=MONO, size=10, bold=True, color=BLUE_INK)
    add_textbox(s, Inches(5.6), Inches(4.5), Inches(2.8), Inches(1.4),
                "Jumlah bahan/upah untuk 1 unit output. Contoh: 371 kg semen + 0,499 m³ pasir + 0,776 m³ kerikil + 215 L air = 1 m³ beton K-225.",
                font=SANS, size=11, color=INK2)

    slide_caption(s, "Rumus AHSP",
                  "OHP 10% menutup biaya tak langsung kontraktor + laba. Harga satuan ini dipakai di RAB.",
                  cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK)


def slide_ahsp_example(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 5, "Tahap 1 · Harga Satuan", cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK)
    slide_title(s, "Contoh Nyata: 1 m³ Beton Struktur K-225",
                "Pecahkan resep menjadi angka — koefisien dari SNI 7394, harga dari HSPK rata-rata Jawa.")

    # table
    rows = [
        ("A · BAHAN (MATERIAL)", None, None, None, None, "head", BLUE_SFT, BLUE_INK),
        ("Semen Portland", "371", "kg", "Rp 1.500", "Rp 556.500", "data", None, None),
        ("Pasir beton", "0,499", "m³", "Rp 285.000", "Rp 142.215", "data", None, None),
        ("Kerikil / split", "0,776", "m³", "Rp 320.000", "Rp 248.320", "data", None, None),
        ("Air", "215", "L", "Rp 50", "Rp 10.750", "data", None, None),
        ("B · UPAH (LABOUR)", None, None, None, None, "head", TEAL_SFT, TEAL_INK),
        ("Pekerja", "1,650", "OH", "Rp 110.000", "Rp 181.500", "data", None, None),
        ("Tukang batu", "0,275", "OH", "Rp 150.000", "Rp 41.250", "data", None, None),
        ("Kepala tukang", "0,028", "OH", "Rp 170.000", "Rp 4.760", "data", None, None),
        ("Mandor", "0,083", "OH", "Rp 180.000", "Rp 14.940", "data", None, None),
        ("Subtotal langsung", "", "", "", "Rp 1.200.235", "sub", None, None),
        ("Overhead & Profit (10%)", "", "", "", "+ Rp 120.024", "data", None, None),
        ("HARGA SATUAN per m³", "", "", "", "Rp 1.320.259", "total", BLUE, WHITE),
    ]

    table_left = Inches(0.6); table_top = Inches(2.4); table_w = Inches(12.1)
    row_h = Inches(0.30)
    # column widths
    cw = [Inches(4.0), Inches(1.4), Inches(0.8), Inches(2.4), Inches(3.5)]
    # header row
    hdr_top = table_top
    hdr_rect = add_rect(s, table_left, hdr_top, table_w, Inches(0.32),
                        fill=SURF2, line=LINE, line_w=0.5)
    headers = ["Komponen", "Koef.", "Sat.", "Harga Dasar", "Jumlah"]
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
    cx = table_left
    for i, (h, a) in enumerate(zip(headers, aligns)):
        add_textbox(s, cx + Inches(0.1), hdr_top + Inches(0.04), cw[i] - Inches(0.2), Inches(0.28),
                    h.upper(), font=MONO, size=9, bold=True, color=MUTED, align=a)
        cx += cw[i]

    # body rows
    y = hdr_top + Inches(0.32)
    for row in rows:
        text, coef, unit, price, total, kind, bg, ink = row
        if kind == "head":
            add_rect(s, table_left, y, table_w, Inches(0.30),
                     fill=bg, line=None)
            add_textbox(s, table_left + Inches(0.1), y + Inches(0.04), table_w, Inches(0.28),
                        text, font=MONO, size=9, bold=True, color=ink)
        elif kind == "total":
            add_rect(s, table_left, y, table_w, Inches(0.36),
                     fill=bg, line=None)
            # label
            add_textbox(s, table_left + Inches(0.1), y + Inches(0.06), Inches(8.0), Inches(0.3),
                        text, font=SANS, size=12, bold=True, color=ink)
            # value at right
            add_textbox(s, table_left + sum(cw[:4]), y + Inches(0.06), cw[4] - Inches(0.1), Inches(0.3),
                        total, font=MONO, size=13, bold=True, color=ink, align=PP_ALIGN.RIGHT)
            y += Inches(0.06)  # bump for taller row
        elif kind == "sub":
            add_rect(s, table_left, y, table_w, Inches(0.32),
                     fill=SURF2, line=INK, line_w=0.7)
            add_textbox(s, table_left + Inches(0.1), y + Inches(0.04), Inches(8.0), Inches(0.28),
                        text, font=SANS, size=11, bold=True, color=INK)
            add_textbox(s, table_left + sum(cw[:4]), y + Inches(0.04), cw[4] - Inches(0.1), Inches(0.28),
                        total, font=MONO, size=11, bold=True, color=INK, align=PP_ALIGN.RIGHT)
        else:
            # normal data row
            add_rect(s, table_left, y, table_w, Inches(0.28),
                     fill=SURF, line=LINE, line_w=0.4)
            vals = [text, coef, unit, price, total]
            cx2 = table_left
            for i, (v, a) in enumerate(zip(vals, aligns)):
                if v:
                    font = MONO if i > 0 else SANS
                    add_textbox(s, cx2 + Inches(0.1), y + Inches(0.03), cw[i] - Inches(0.2), Inches(0.25),
                                v, font=font, size=10, color=INK if i == 0 else INK2, align=a)
                cx2 += cw[i]
        y += Inches(0.30 if kind != "total" else 0.36)

    slide_caption(s, "Studi Kasus",
                  "Angka ini dipakai di seluruh RAB untuk setiap m³ beton struktur. Volume 32 m³ × Rp 1.320.259 = Rp 42,2 jt.",
                  cat_color=BLUE, cat_soft=BLUE_SFT, cat_ink=BLUE_INK)


def slide_rab_formula(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 6, "Tahap 2 · RAB", cat_color=TEAL, cat_soft=TEAL_SFT, cat_ink=TEAL_INK)
    slide_title(s, "Cara Membuat RAB: Kalikan dan Jumlahkan",
                "RAB = Rencana Anggaran Biaya. Volume × Harga Satuan, lalu jumlahkan semua item.")

    # 3-pill equation at top
    eq_top = Inches(2.5)
    boxes = [
        ("DARI REVIT BIM", "Volume",       TEAL_SFT, TEAL,  TEAL_INK, False),
        ("DARI AHSP",      "Harga Satuan", BLUE_SFT, BLUE,  BLUE_INK, False),
        ("HASIL PER ITEM", "Total Item",   TEAL,      None,  WHITE, True),
    ]
    bx = Inches(1.8); bw = Inches(2.6); bh = Inches(0.9)
    cx = bx
    for i, (eb, t, fill, line, ink, solid) in enumerate(boxes):
        add_rect(s, cx, eq_top, bw, bh,
                 fill=fill, line=line, line_w=1.5, radius=True)
        add_textbox(s, cx, eq_top + Inches(0.15), bw, Inches(0.3),
                    eb, font=MONO, size=10, bold=True, color=ink if not solid else WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, cx, eq_top + Inches(0.5), bw, Inches(0.4),
                    t, font=SANS, size=18, bold=True, color=INK if not solid else WHITE, align=PP_ALIGN.CENTER)
        cx += bw
        if i < 2:
            sym = "×" if i == 0 else "="
            add_textbox(s, cx - Inches(0.05), eq_top + Inches(0.3), Inches(0.5), Inches(0.5),
                        sym, font=MONO, size=24, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
            cx += Inches(0.4)

    # 3 sample rows
    items = [
        ("D.1 Beton struktur",      "32,40 m³",  "Rp 1.320.259", "Rp 42.776.392"),
        ("E.1 Dinding bata merah",  "218,00 m²", "Rp 160.000",   "Rp 34.880.000"),
        ("H.1 Pengecatan interior", "434,00 m²", "Rp 45.000",    "Rp 19.530.000"),
    ]
    rl = Inches(0.6); rt = Inches(4.0); rw = Inches(12.1); rh = Inches(0.34)
    for i, (lbl, vol, hs, tot) in enumerate(items):
        ty = rt + (rh + Inches(0.05)) * i
        add_rect(s, rl, ty, rw, rh, fill=SURF2, line=LINE, line_w=0.5)
        add_textbox(s, rl + Inches(0.15), ty + Inches(0.04), Inches(4), Inches(0.3),
                    lbl, font=SANS, size=11, color=INK)
        add_textbox(s, rl + Inches(4.3), ty + Inches(0.04), Inches(1.8), Inches(0.3),
                    vol, font=MONO, size=11, color=INK, align=PP_ALIGN.RIGHT)
        add_textbox(s, rl + Inches(6.3), ty + Inches(0.04), Inches(0.3), Inches(0.3),
                    "×", font=MONO, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_textbox(s, rl + Inches(6.7), ty + Inches(0.04), Inches(2.2), Inches(0.3),
                    hs, font=MONO, size=11, color=INK, align=PP_ALIGN.RIGHT)
        add_textbox(s, rl + Inches(9.0), ty + Inches(0.04), Inches(0.3), Inches(0.3),
                    "=", font=MONO, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_textbox(s, rl + Inches(9.4), ty + Inches(0.04), Inches(2.6), Inches(0.3),
                    tot, font=MONO, size=11, bold=True, color=INK, align=PP_ALIGN.RIGHT)

    # sum row
    sy = rt + (rh + Inches(0.05)) * 3 + Inches(0.1)
    add_rect(s, rl, sy, rw, Inches(0.5), fill=TEAL, line=None, radius=True)
    add_textbox(s, rl + Inches(0.15), sy + Inches(0.12), Inches(8), Inches(0.3),
                "Σ semua item + PPN 11% → JUMLAH BIAYA PROYEK",
                font=SANS, size=12, bold=True, color=WHITE)
    add_textbox(s, rl + Inches(9.0), sy + Inches(0.12), Inches(3.0), Inches(0.3),
                "Rp 532.755.000",
                font=MONO, size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    slide_caption(s, "Cara Kerja RAB",
                  "Setiap baris = satu item WBS dengan volume dari Revit, harga dari AHSP. Total + PPN = anggaran proyek.",
                  cat_color=TEAL, cat_soft=TEAL_SFT, cat_ink=TEAL_INK)


def slide_rab_markup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 7, "Tahap 2 · RAB", cat_color=TEAL, cat_soft=TEAL_SFT, cat_ink=TEAL_INK)
    slide_title(s, "Setelah Item Dijumlahkan: Naikkan ke Total Proyek",
                "Subtotal 'biaya langsung' belum termasuk pajak. Tambahkan markup tingkat proyek.")

    # stair-step: 3 boxes
    steps = [
        ("LANGKAH 1", "Jumlah Biaya Langsung", "Σ (Volume × Harga Satuan)",   "Rp 480.000.000", TEAL_SFT, TEAL,  TEAL_INK,  False),
        ("LANGKAH 2", "+ PPN 11%",              "Pajak Pertambahan Nilai",      "+ Rp 52.800.000", TEAL_SFT, TEAL,  TEAL_INK,  False),
        ("TOTAL",     "TOTAL PROYEK",           "Yang dibayar pemilik ke kontraktor", "Rp 532,8 jt", TEAL, None, WHITE, True),
    ]
    bx = Inches(0.8); by = [Inches(4.2), Inches(3.5), Inches(2.9)]
    bw = Inches(3.5); bh = Inches(1.5)
    for i, (eb, name, sub, val, fill, line, ink, solid) in enumerate(steps):
        left = bx + Inches(3.9) * i
        top = by[i]
        add_rect(s, left, top, bw, bh, fill=fill, line=line, line_w=1.5, radius=True)
        add_textbox(s, left + Inches(0.2), top + Inches(0.15), bw, Inches(0.3),
                    eb, font=MONO, size=10, bold=True, color=ink if not solid else BLUE_SFT)
        add_textbox(s, left + Inches(0.2), top + Inches(0.45), bw, Inches(0.4),
                    name, font=SANS, size=16, bold=True, color=INK if not solid else WHITE)
        add_textbox(s, left + Inches(0.2), top + Inches(0.8), bw, Inches(0.3),
                    sub, font=SANS, size=10, color=MUTED if not solid else TEAL_SFT)
        add_textbox(s, left + Inches(0.2), top + Inches(1.05), bw, Inches(0.35),
                    val, font=MONO, size=15 if solid else 13, bold=True,
                    color=INK if not solid else WHITE)

    # arrows between
    for i in range(2):
        x1 = bx + Inches(3.9) * i + bw
        y1 = by[i] + bh / 2
        x2 = bx + Inches(3.9) * (i + 1)
        y2 = by[i + 1] + bh / 2
        add_arrow(s, x1, y1, x2, y2, color=TEAL, width=2)
        # label
        mark = "+11%" if i == 0 else "="
        add_textbox(s, (x1 + x2) / 2 - Inches(0.4), (y1 + y2) / 2 - Inches(0.35),
                    Inches(0.8), Inches(0.3),
                    mark, font=MONO, size=10, bold=True, color=TEAL_INK, align=PP_ALIGN.CENTER)

    # TW footnote
    add_textbox(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
                "📌  Standar Taiwan: bukan 1 baris PPN, tapi 5 lapis markup (管理 7% + 利潤 5% + 品管 1% + 環安衛 1,5% + 營業稅 5%).",
                font=SANS, size=11, color=MUTED)

    slide_caption(s, "Markup Proyek",
                  "Indonesia menyisipkan OHP 10% di tiap harga satuan via AHSP — di tingkat proyek hanya PPN 11%.",
                  cat_color=TEAL, cat_soft=TEAL_SFT, cat_ink=TEAL_INK)


def slide_sched_duration(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 8, "Tahap 3 · Jadwal", cat_color=AMBER, cat_soft=AMBER_SFT, cat_ink=AMBER_INK)
    slide_title(s, "Cara Menghitung Durasi: Volume Dibagi Kecepatan",
                "Bukan tebakan. Setiap pekerjaan punya angka produktivitas baku per hari per regu.")

    # big formula
    add_textbox(s, Inches(1.2), Inches(2.6), Inches(3.0), Inches(0.6),
                "Durasi (hari) =", font=SANS, size=22, bold=True, color=INK)
    # numerator
    add_textbox(s, Inches(4.4), Inches(2.5), Inches(4.5), Inches(0.4),
                "Volume Pekerjaan", font=SANS, size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(4.4), Inches(2.85), Inches(4.5), Inches(0.3),
                "(m³, m², kg)", font=MONO, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # bar
    bar = s.shapes.add_connector(1, Inches(4.4), Inches(3.2), Inches(8.9), Inches(3.2))
    bar.line.color.rgb = INK; bar.line.width = Pt(2)
    # denominator
    add_textbox(s, Inches(4.4), Inches(3.25), Inches(4.5), Inches(0.4),
                "Produktivitas × Jumlah Regu", font=SANS, size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(4.4), Inches(3.60), Inches(4.5), Inches(0.3),
                "(m³/hari/regu)", font=MONO, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # worked example boxes
    add_textbox(s, Inches(0.6), Inches(4.4), Inches(11.5), Inches(0.3),
                "📐 CONTOH NYATA — Beton Struktur D.1",
                font=MONO, size=11, bold=True, color=AMBER_INK)
    ex_top = Inches(4.7)
    ex_items = [
        ("VOLUME (Revit)", "32 m³",            AMBER_SFT, AMBER),
        ("PRODUKTIVITAS",  "1,2 m³/hari",      AMBER_SFT, AMBER),
        ("DURASI",         "≈ 27 hari = 6 minggu", AMBER, None),
    ]
    bw = Inches(3.8); bh = Inches(0.9); gap = Inches(0.15); left = Inches(0.6)
    for i, (eb, val, fill, line) in enumerate(ex_items):
        lx = left + (bw + gap) * i
        add_rect(s, lx, ex_top, bw, bh, fill=fill, line=line, line_w=1.2, radius=True)
        ink = WHITE if line is None else AMBER_INK
        bg_ink = AMBER_SFT if line is None else AMBER_INK
        add_textbox(s, lx + Inches(0.2), ex_top + Inches(0.13), bw, Inches(0.3),
                    eb, font=MONO, size=10, bold=True, color=bg_ink)
        add_textbox(s, lx + Inches(0.2), ex_top + Inches(0.4), bw, Inches(0.45),
                    val, font=MONO, size=18, bold=True, color=WHITE if line is None else INK)

    # productivity table
    add_textbox(s, Inches(0.6), Inches(5.85), Inches(11.5), Inches(0.3),
                "📊 PRODUKTIVITAS BAKU TIAP PEKERJAAN:",
                font=MONO, size=11, bold=True, color=AMBER_INK)
    items = [
        ("Cor beton",   "1,2 m³/hari"),
        ("Pasang bata", "10 m²/hari"),
        ("Plester",     "12 m²/hari"),
        ("Pengecatan",  "25 m²/hari"),
    ]
    py = Inches(6.15)
    for i, (n, v) in enumerate(items):
        lx = Inches(0.6) + Inches(3.0) * i
        add_textbox(s, lx, py, Inches(1.5), Inches(0.3),
                    n + ":", font=SANS, size=11, color=INK)
        add_textbox(s, lx + Inches(1.4), py, Inches(1.5), Inches(0.3),
                    v, font=MONO, size=11, color=AMBER_INK)

    slide_caption(s, "Rumus Durasi",
                  "Curing beton menambah waktu kalender meski tidak ada 'kerja'. 27 hari dibulatkan ke 6 minggu.",
                  cat_color=AMBER, cat_soft=AMBER_SFT, cat_ink=AMBER_INK)


def slide_sched_predecessor(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 9, "Tahap 3 · Jadwal", cat_color=AMBER, cat_soft=AMBER_SFT, cat_ink=AMBER_INK)
    slide_title(s, "Urutan Pekerjaan: Apa Harus Selesai Dulu?",
                "Predecessor adalah ketergantungan fisik-logis (kolom butuh pondasi dulu).")

    # 8 nodes A→H in a horizontal flow
    nodes = [
        ("A", "Persiapan", False),
        ("B", "Galian", False),
        ("C", "Pondasi", False),
        ("D", "Struktur", True),    # filled
        ("E", "Arsitektur", False),
        ("F+G", "Atap & MEP", False),
        ("H", "Finishing", False),
    ]
    bw = Inches(1.5); bh = Inches(1.0); gap = Inches(0.18)
    total = bw * len(nodes) + gap * (len(nodes) - 1)
    start = (SLIDE_W - total) / 2
    top = Inches(3.2)

    for i, (letter, name, highlight) in enumerate(nodes):
        left = start + (bw + gap) * i
        fill = AMBER if highlight else AMBER_SFT
        line = None if highlight else AMBER
        add_rect(s, left, top, bw, bh, fill=fill, line=line, line_w=1.5, radius=True)
        ink = WHITE if highlight else INK
        add_textbox(s, left, top + Inches(0.15), bw, Inches(0.45),
                    letter, font=SANS, size=18, bold=True, color=ink, align=PP_ALIGN.CENTER)
        add_textbox(s, left, top + Inches(0.6), bw, Inches(0.35),
                    name, font=SANS, size=11, color=ink, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            x1 = left + bw
            x2 = start + (bw + gap) * (i + 1)
            ym = top + bh / 2
            add_arrow(s, x1, ym, x2, ym, color=AMBER, width=1.5)
            add_textbox(s, (x1 + x2) / 2 - Inches(0.2), ym - Inches(0.35),
                        Inches(0.4), Inches(0.25),
                        "FS", font=MONO, size=8, bold=True, color=AMBER_INK, align=PP_ALIGN.CENTER)

    # critical path label above D
    add_textbox(s, start + (bw + gap) * 3, top - Inches(0.45), bw, Inches(0.3),
                "⚡ LINTASAN KRITIS", font=MONO, size=11, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)

    # legend bottom
    add_textbox(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.3),
                "EMPAT TIPE HUBUNGAN:", font=MONO, size=11, bold=True, color=AMBER_INK)
    legends = [
        ("FS Finish-to-Start — paling umum, A selesai → B mulai",     Inches(0.8), Inches(5.4)),
        ("SS Start-to-Start — A & B mulai bersamaan",                  Inches(0.8), Inches(5.7)),
        ("FF Finish-to-Finish — A & B selesai bersamaan",              Inches(6.8), Inches(5.4)),
        ("SF Start-to-Finish — jarang dipakai",                        Inches(6.8), Inches(5.7)),
    ]
    for txt, lx, ly in legends:
        add_textbox(s, lx, ly, Inches(5.5), Inches(0.3),
                    txt, font=SANS, size=11, color=INK2)

    slide_caption(s, "Predecessor",
                  "Lintasan kritis menentukan total durasi. Bila satu pekerjaan di lintasan kritis mundur, seluruh proyek mundur.",
                  cat_color=AMBER, cat_soft=AMBER_SFT, cat_ink=AMBER_INK)


def slide_sched_gantt(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 10, "Tahap 3 · Jadwal", cat_color=AMBER, cat_soft=AMBER_SFT, cat_ink=AMBER_INK)
    slide_title(s, "Hasil Akhir: Bagan Gantt",
                "Sumbu X = minggu. Sumbu Y = pekerjaan. Panjang bar = durasi. Posisi bar = kapan dikerjakan.")

    # Gantt area
    g_left = Inches(2.3); g_top = Inches(2.7); g_w = Inches(10.5); g_h = Inches(3.3)
    weeks = 18
    col_w = g_w / weeks

    # week labels
    for w in range(1, weeks + 1):
        add_textbox(s, g_left + col_w * (w - 1), g_top - Inches(0.3), col_w, Inches(0.25),
                    f"W{w}", font=MONO, size=8, color=MUTED, align=PP_ALIGN.CENTER)

    # vertical gridlines
    for w in range(weeks + 1):
        x = g_left + col_w * w
        line = s.shapes.add_connector(1, x, g_top, x, g_top + g_h)
        line.line.color.rgb = LINE; line.line.width = Pt(0.3)

    # rows
    rows = [
        ("A · Persiapan",    1, 2,  BLUE),
        ("B · Tanah",        2, 4,  TEAL),
        ("C · Pondasi",      3, 4,  AMBER),
        ("D · Struktur",     4, 9,  ROSE),
        ("E · Arsitektur",   8, 15, GREEN),
        ("F · Atap",         9, 12, VIOLET),
        ("G · MEP",          10, 16, BLUE),
        ("H · Finishing",    14, 18, AMBER),
    ]
    row_h = g_h / (len(rows) + 1)
    for i, (lbl, sw, ew, color) in enumerate(rows):
        row_y = g_top + row_h * i + row_h * 0.15
        # label on left
        add_textbox(s, Inches(0.5), row_y - Inches(0.05), Inches(1.7), Inches(0.3),
                    lbl, font=SANS, size=10, bold=True, color=INK, align=PP_ALIGN.RIGHT)
        # bar
        bx = g_left + col_w * (sw - 1)
        bw = col_w * (ew - sw + 1)
        bar = add_rect(s, bx, row_y, bw, row_h * 0.65,
                       fill=color, line=None, radius=True)

    # TOTAL bar (the union)
    total_y = g_top + row_h * len(rows) + Inches(0.15)
    add_textbox(s, Inches(0.5), total_y - Inches(0.02), Inches(1.7), Inches(0.3),
                "TOTAL", font=SANS, size=10, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    add_rect(s, g_left, total_y, g_w, Inches(0.32), fill=INK, line=None, radius=True)
    add_textbox(s, g_left, total_y + Inches(0.04), g_w, Inches(0.25),
                "18 minggu (≈ 126 hari kalender)",
                font=SANS, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_caption(s, "Gantt Chart",
                  "Warna bar = grup WBS. Bar tumpang tindih = pekerjaan paralel oleh regu berbeda.",
                  cat_color=AMBER, cat_soft=AMBER_SFT, cat_ink=AMBER_INK)


def slide_scurve_weight(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 11, "Tahap 4 · Kurva S", cat_color=ROSE, cat_soft=ROSE_SFT, cat_ink=ROSE_INK)
    slide_title(s, "Langkah 1 — Hitung 'Bobot' Tiap Pekerjaan",
                "Bobot = porsi biaya item terhadap total proyek. Total semua bobot = 100%.")

    # Big formula at top
    add_textbox(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(0.4),
                "Bobot% = (Total item ÷ Total proyek) × 100",
                font=MONO, size=15, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)

    # Horizontal stacked bar showing weights
    bar_left = Inches(0.6); bar_top = Inches(3.2); bar_w = Inches(12.1); bar_h = Inches(0.6)
    # data: A 4, B 6, C 14, D 28, E 22, F 8, G 10, H 8
    data = [
        ("A · Persiapan",   4,  BLUE),
        ("B · Tanah",       6,  TEAL),
        ("C · Pondasi",     14, AMBER),
        ("D · Struktur",    28, ROSE),
        ("E · Arsitektur",  22, GREEN),
        ("F · Atap",        8,  VIOLET),
        ("G · MEP",         10, BLUE),
        ("H · Finishing",   8,  AMBER),
    ]
    cx = bar_left
    for i, (lbl, pct, color) in enumerate(data):
        seg_w = bar_w * pct / 100
        add_rect(s, cx, bar_top, seg_w, bar_h, fill=color, line=None)
        # label on top of bar
        if pct >= 6:
            add_textbox(s, cx, bar_top + Inches(0.16), seg_w, Inches(0.28),
                        f"{pct}%", font=MONO, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += seg_w

    # Legend grid below
    lg_top = bar_top + bar_h + Inches(0.4)
    cols = 4
    cw = Inches(3.1)
    for i, (lbl, pct, color) in enumerate(data):
        col = i % cols; row = i // cols
        lx = bar_left + cw * col
        ly = lg_top + Inches(0.45) * row
        # color swatch
        sw = add_rect(s, lx, ly + Inches(0.04), Inches(0.16), Inches(0.16),
                      fill=color, line=None)
        add_textbox(s, lx + Inches(0.22), ly, cw - Inches(1.0), Inches(0.3),
                    lbl, font=SANS, size=10, color=INK)
        add_textbox(s, lx + cw - Inches(0.7), ly, Inches(0.6), Inches(0.3),
                    f"{pct}%", font=MONO, size=11, bold=True, color=INK, align=PP_ALIGN.RIGHT)

    slide_caption(s, "Bobot",
                  "Struktur (D) + Arsitektur (E) = 50% biaya proyek. Itu sebabnya keduanya jadi inti fase 'curam' kurva S.",
                  cat_color=ROSE, cat_soft=ROSE_SFT, cat_ink=ROSE_INK)


def slide_scurve_distribution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 12, "Tahap 4 · Kurva S", cat_color=ROSE, cat_soft=ROSE_SFT, cat_ink=ROSE_INK)
    slide_title(s, "Langkah 2 — Sebarkan Bobot ke Sepanjang Durasinya",
                "Bobot 8,8% dengan durasi 6 minggu → 1,47% per minggu. Lakukan untuk semua item.")

    # formula
    add_textbox(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(0.4),
                "Bobot / minggu = Bobot item ÷ Durasi",
                font=MONO, size=14, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)

    # 3-step example for D.1
    add_textbox(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(0.3),
                "CONTOH UNTUK D.1 BETON STRUKTUR:",
                font=MONO, size=11, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)
    steps = [
        ("BOBOT TOTAL", "28%",       ROSE_SFT, ROSE),
        ("DIBAGI",      "6 minggu",  ROSE_SFT, ROSE),
        ("PER MINGGU",  "≈ 4,7 %",   ROSE, None),
    ]
    bw = Inches(3.8); bh = Inches(0.9); gap = Inches(0.15); left = Inches(0.6)
    eq_y = Inches(3.7)
    for i, (eb, val, fill, line) in enumerate(steps):
        lx = left + (bw + gap) * i
        add_rect(s, lx, eq_y, bw, bh, fill=fill, line=line, line_w=1.2, radius=True)
        is_solid = line is None
        ink = ROSE_INK if not is_solid else ROSE_SFT
        add_textbox(s, lx + Inches(0.2), eq_y + Inches(0.12), bw, Inches(0.3),
                    eb, font=MONO, size=10, bold=True, color=ink)
        add_textbox(s, lx + Inches(0.2), eq_y + Inches(0.4), bw, Inches(0.45),
                    val, font=MONO, size=20, bold=True, color=INK if not is_solid else WHITE)
        if i < 2:
            sym = "÷" if i == 0 else "="
            add_textbox(s, lx + bw + Inches(0.02), eq_y + Inches(0.25), Inches(0.4), Inches(0.4),
                        sym, font=MONO, size=22, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # mini matrix below
    add_textbox(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.3),
                "JADINYA, MINGGU 4-9 (6 KOLOM) MASING-MASING DAPAT 4,7%:",
                font=MONO, size=10, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)

    # mini week strip for D.1
    weeks = 18
    strip_left = Inches(2.0); strip_top = Inches(5.4); strip_w = Inches(9.5); strip_h = Inches(0.5)
    col_w = strip_w / weeks
    # labels
    for w in range(1, weeks + 1):
        add_textbox(s, strip_left + col_w * (w - 1), strip_top - Inches(0.25), col_w, Inches(0.2),
                    str(w), font=MONO, size=8, color=FAINT, align=PP_ALIGN.CENTER)
    # cells — D.1 has weight in weeks 4-9
    for w in range(1, weeks + 1):
        active = 4 <= w <= 9
        cell = add_rect(s, strip_left + col_w * (w - 1), strip_top, col_w, strip_h,
                        fill=ROSE if active else SURF, line=LINE, line_w=0.4)
        add_textbox(s, strip_left + col_w * (w - 1), strip_top + Inches(0.12), col_w, Inches(0.3),
                    "4,7" if active else "·", font=MONO, size=10, bold=active,
                    color=WHITE if active else FAINT, align=PP_ALIGN.CENTER)

    # row label
    add_textbox(s, Inches(0.5), strip_top + Inches(0.1), Inches(1.4), Inches(0.3),
                "D.1 Beton:", font=SANS, size=11, bold=True, color=ROSE_INK, align=PP_ALIGN.RIGHT)

    slide_caption(s, "Distribusi",
                  "Ulangi untuk semua item. Jumlah satu kolom (semua item aktif di minggu itu) = 'Rencana/Minggu'.",
                  cat_color=ROSE, cat_soft=ROSE_SFT, cat_ink=ROSE_INK)


def slide_scurve_curve(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 13, "Tahap 4 · Kurva S", cat_color=ROSE, cat_soft=ROSE_SFT, cat_ink=ROSE_INK)
    slide_title(s, "Langkah 3 — Akumulasi: Inilah Kurva S",
                "Jumlahkan 'Rencana per Minggu' dari minggu 1 → seterusnya. Hasilnya membentuk huruf S.")

    # chart frame
    chart_left = Inches(1.5); chart_top = Inches(2.7); chart_w = Inches(10.5); chart_h = Inches(3.4)

    # Y axis: 0-100% in 4 grid lines
    for i, val in enumerate([0, 25, 50, 75, 100]):
        y = chart_top + chart_h - (chart_h * val / 100)
        # gridline
        line = s.shapes.add_connector(1, chart_left, y, chart_left + chart_w, y)
        line.line.color.rgb = INK if val == 0 else LINE
        line.line.width = Pt(1.0 if val == 0 else 0.4)
        # label
        add_textbox(s, chart_left - Inches(0.5), y - Inches(0.13), Inches(0.4), Inches(0.25),
                    str(val), font=MONO, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    # X axis labels (1-18)
    weeks = 18
    col_w = chart_w / (weeks - 1)
    for w in range(1, weeks + 1):
        x = chart_left + col_w * (w - 1)
        add_textbox(s, x - Inches(0.15), chart_top + chart_h + Inches(0.05), Inches(0.3), Inches(0.2),
                    str(w), font=MONO, size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # cumulative S-curve (sample data points)
    cum = [2, 6, 15, 28.7, 33.4, 38.1, 42.8, 50.3, 59.8, 64.6, 69.4, 74.2, 77, 81.4, 85.8, 87.4, 89, 100]
    # draw freeform / connected segments
    pts = []
    for w, v in enumerate(cum):
        x = chart_left + col_w * w
        y = chart_top + chart_h - (chart_h * v / 100)
        pts.append((x, y))
    # connect with line segments
    for i in range(len(pts) - 1):
        seg = s.shapes.add_connector(1, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1])
        seg.line.color.rgb = ROSE; seg.line.width = Pt(2.5)
    # add dots
    for (x, y) in pts:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.05), y - Inches(0.05), Inches(0.1), Inches(0.1))
        set_fill(dot, ROSE); no_line(dot)

    # phase annotations below
    add_textbox(s, chart_left, chart_top + chart_h + Inches(0.45), Inches(3.0), Inches(0.3),
                "FASE 1 (landai)", font=SANS, size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, chart_left, chart_top + chart_h + Inches(0.7), Inches(3.0), Inches(0.25),
                "Persiapan + Galian", font=MONO, size=9, color=FAINT, align=PP_ALIGN.CENTER)

    add_textbox(s, chart_left + Inches(3.5), chart_top + chart_h + Inches(0.45), Inches(3.5), Inches(0.3),
                "FASE 2 (curam — kurva S!)", font=SANS, size=11, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)
    add_textbox(s, chart_left + Inches(3.5), chart_top + chart_h + Inches(0.7), Inches(3.5), Inches(0.25),
                "Pondasi + Struktur + Arsitektur", font=MONO, size=9, color=FAINT, align=PP_ALIGN.CENTER)

    add_textbox(s, chart_left + Inches(7.5), chart_top + chart_h + Inches(0.45), Inches(3.0), Inches(0.3),
                "FASE 3 (landai)", font=SANS, size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, chart_left + Inches(7.5), chart_top + chart_h + Inches(0.7), Inches(3.0), Inches(0.25),
                "Finishing + Cleanup", font=MONO, size=9, color=FAINT, align=PP_ALIGN.CENTER)

    # Y axis title
    add_textbox(s, Inches(0.4), Inches(4.2), Inches(0.8), Inches(0.4),
                "PROGRES\nKUMULATIF (%)", font=MONO, size=9, bold=True, color=ROSE_INK, align=PP_ALIGN.CENTER)

    slide_caption(s, "Kurva S",
                  "Tiga fase (landai → curam → landai) memberi bentuk huruf S. Kurva ini jadi baseline progres.",
                  cat_color=ROSE, cat_soft=ROSE_SFT, cat_ink=ROSE_INK)


def slide_recap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, 14, "Recap", cat_color=GREEN, cat_soft=GREEN_SFT, cat_ink=GREEN_INK)
    slide_title(s, "Semua Saling Berkaitan",
                "Ubah satu tahap, semua hilirnya ikut berubah. Itulah kekuatan 5D BIM.")

    # center node (Revit BIM)
    cx = Inches(5.4); cy = Inches(3.5); cw = Inches(2.5); ch = Inches(1.4)
    add_rect(s, cx, cy, cw, ch, fill=GREEN, line=None, radius=True)
    add_textbox(s, cx, cy + Inches(0.2), cw, Inches(0.3),
                "MODEL 3D", font=MONO, size=11, bold=True, color=GREEN_SFT, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(0.55), cw, Inches(0.5),
                "Revit BIM", font=SANS, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 4 corner output boxes
    corners = [
        # (eyebrow, name, formula, fill, line, ink, left, top)
        ("TAHAP 1", "Harga Satuan", "Σ koefisien × harga + OHP",  BLUE_SFT, BLUE, BLUE_INK,   Inches(0.6),  Inches(2.0)),
        ("TAHAP 2", "RAB",          "Volume × Harga + Markup",     TEAL_SFT, TEAL, TEAL_INK,   Inches(9.5),  Inches(2.0)),
        ("TAHAP 3", "Jadwal · Gantt","Volume ÷ Produktivitas",     AMBER_SFT, AMBER, AMBER_INK, Inches(0.6),  Inches(5.5)),
        ("TAHAP 4", "Kurva S",      "Σ bobot/minggu kumulatif",    ROSE_SFT, ROSE, ROSE_INK,   Inches(9.5),  Inches(5.5)),
    ]
    bw = Inches(3.2); bh = Inches(1.1)
    for (eb, name, formula, fill, line, ink, lx, ly) in corners:
        add_rect(s, lx, ly, bw, bh, fill=fill, line=line, line_w=1.5, radius=True)
        add_textbox(s, lx + Inches(0.2), ly + Inches(0.12), bw, Inches(0.3),
                    eb, font=MONO, size=10, bold=True, color=ink)
        add_textbox(s, lx + Inches(0.2), ly + Inches(0.4), bw, Inches(0.4),
                    name, font=SANS, size=15, bold=True, color=INK)
        add_textbox(s, lx + Inches(0.2), ly + Inches(0.78), bw, Inches(0.3),
                    formula, font=MONO, size=9, color=MUTED)
        # arrow toward center
        mid_x = lx + bw / 2; mid_y = ly + bh / 2
        center_x = cx + cw / 2; center_y = cy + ch / 2
        # endpoint closest to center
        if lx < cx:  # left side
            ex = lx + bw; ey = ly + bh / 2
        else:        # right side
            ex = lx;     ey = ly + bh / 2
        add_arrow(s, ex, ey, center_x + (-cw / 2 if lx < cx else cw / 2), center_y, color=GREEN, width=1.5)

    # center caption
    add_textbox(s, cx, cy + ch + Inches(0.1), cw, Inches(0.3),
                "satu model, empat output, satu kebenaran",
                font=SANS, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    slide_caption(s, "Penutup",
                  "Itulah inti 5D BIM Construction Management: model 3D adalah jantung; setiap perubahan menjalar otomatis. Terima kasih.",
                  cat_color=GREEN, cat_soft=GREEN_SFT, cat_ink=GREEN_INK)


# ─── main ──────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_pipeline(prs)
    slide_ahsp_sources(prs)
    slide_ahsp_formula(prs)
    slide_ahsp_example(prs)
    slide_rab_formula(prs)
    slide_rab_markup(prs)
    slide_sched_duration(prs)
    slide_sched_predecessor(prs)
    slide_sched_gantt(prs)
    slide_scurve_weight(prs)
    slide_scurve_distribution(prs)
    slide_scurve_curve(prs)
    slide_recap(prs)

    out = "Construction Management Process.pptx"
    prs.save(out)
    print(f"[OK] Generated: {out}")
    print(f"     {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
